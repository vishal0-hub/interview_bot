"""
Generate a PowerPoint presentation for the AI Interview Bot project.
Run: py create_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Theme Colors ───
BG_DARK = RGBColor(10, 14, 26)
BG_CARD = RGBColor(17, 24, 39)
ACCENT = RGBColor(99, 102, 241)
ACCENT_LIGHT = RGBColor(139, 92, 246)
WHITE = RGBColor(241, 245, 249)
GRAY = RGBColor(148, 163, 184)
GREEN = RGBColor(16, 185, 129)
YELLOW = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Helpers ───
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_shape_bg(slide, left, top, width, height, color=BG_CARD, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.05
    return shape

def add_bullet_list(slide, left, top, width, height, items, size=16, color=GRAY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

# ════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide, BG_DARK)

# Decorative accent circle
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(6), Inches(6))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(30, 27, 75)
shape.line.fill.background()

shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(4), Inches(4), Inches(4))
shape2.fill.solid()
shape2.fill.fore_color.rgb = RGBColor(30, 27, 75)
shape2.line.fill.background()

add_text(slide, 1.5, 0.8, 5, 0.6, '🤖', size=48, alignment=PP_ALIGN.LEFT)
add_text(slide, 1.5, 1.8, 10, 1.2, 'AI-Powered Technical Interview Bot', size=44, bold=True, color=WHITE)
add_accent_line(slide, 1.5, 3.2, 3)
add_text(slide, 1.5, 3.5, 10, 0.8, 'Resume Analysis • Question Generation • Voice Interaction • Performance Report', size=20, color=GRAY)
add_text(slide, 1.5, 5.0, 10, 0.5, 'Built with Django  |  Powered by Groq AI (LLaMA 3.3 70B)', size=16, color=ACCENT_LIGHT)
add_text(slide, 1.5, 6.2, 10, 0.5, 'February 2026', size=14, color=GRAY)

# ════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.8, '📌  Problem Statement', size=36, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.2, 2.5)

add_shape_bg(slide, 0.8, 1.6, 11.5, 1.4)
add_text(slide, 1.2, 1.8, 10.8, 1.2,
    'Traditional technical interviews are time-consuming, inconsistent, and expensive.\n'
    'Companies struggle to screen large volumes of candidates efficiently while maintaining\n'
    'quality and fairness in the evaluation process.',
    size=18, color=GRAY)

problems = [
    '⏱️  Manual screening takes hours per candidate',
    '❌  Inconsistent question quality across interviewers',
    '💰  High cost of technical interviewer time',
    '📊  No standardized scoring or analytics',
    '🔇  No voice-based interaction in existing tools',
]
add_bullet_list(slide, 1.2, 3.5, 10, 3.5, problems, size=18, color=WHITE)

# ════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.8, '💡  Our Solution', size=36, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.2, 2)

add_text(slide, 0.8, 1.6, 11, 0.8,
    'An AI-powered interview bot that automates the entire technical screening process\n'
    '— from resume analysis to final report — using free, state-of-the-art LLMs.',
    size=18, color=GRAY)

# Feature cards
features = [
    ('📄', 'Resume Analysis', 'Upload PDF/DOCX resume.\nAI extracts technical skills\nautomatically.'),
    ('🧠', 'Smart Questions', '10 tailored technical questions\ngenerated based on\ncandidate\'s skill set.'),
    ('🎤', 'Voice Interaction', 'Bot speaks questions aloud.\nCandidate answers via\nmicrophone.'),
    ('📊', 'Analysis Report', 'Detailed scoring, strengths,\nweaknesses, and hiring\nrecommendation.'),
]

for i, (icon, title, desc) in enumerate(features):
    x = 0.8 + i * 3.05
    add_shape_bg(slide, x, 2.8, 2.8, 3.8)
    add_text(slide, x + 0.3, 3.0, 2.2, 0.6, icon, size=36, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.3, 3.7, 2.2, 0.5, title, size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.3, 4.3, 2.2, 2.0, desc, size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 4 — Technology Stack
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.8, '🛠️  Technology Stack', size=36, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.2, 2.5)

stack = [
    ('Backend', ['Django 4.2+ (Python Web Framework)', 'SQLite Database', 'Django ORM for data models']),
    ('AI / LLM', ['Groq API (Free Tier)', 'LLaMA 3.3 70B Versatile Model', 'JSON-structured prompts']),
    ('Resume Parsing', ['PyPDF2 for PDF extraction', 'python-docx for DOCX extraction', 'Custom text parser']),
    ('Frontend', ['HTML5 / CSS3 / JavaScript', 'Web Speech API (TTS + STT)', 'AJAX for real-time Q&A']),
]

for i, (category, items) in enumerate(stack):
    x = 0.8 + (i % 2) * 6.2
    y = 1.6 + (i // 2) * 2.8
    add_shape_bg(slide, x, y, 5.8, 2.4)
    add_text(slide, x + 0.4, y + 0.2, 5, 0.5, category, size=20, bold=True, color=ACCENT)
    for j, item in enumerate(items):
        add_text(slide, x + 0.6, y + 0.8 + j * 0.45, 5, 0.4, f'• {item}', size=14, color=GRAY)

# ════════════════════════════════════════════
# SLIDE 5 — System Architecture / Flow
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.8, '⚙️  System Architecture', size=36, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.2, 2.5)

steps = [
    ('1', 'Upload\nResume', '📄'),
    ('2', 'Parse &\nExtract Text', '📝'),
    ('3', 'AI: Extract\nSkills', '🧠'),
    ('4', 'AI: Generate\nQuestions', '❓'),
    ('5', 'Interactive\nQ&A Session', '🎤'),
    ('6', 'AI: Evaluate\nAnswers', '✅'),
    ('7', 'Generate\nReport', '📊'),
]

for i, (num, label, icon) in enumerate(steps):
    x = 0.5 + i * 1.78
    add_shape_bg(slide, x, 2.0, 1.55, 2.5)
    add_text(slide, x + 0.1, 2.1, 1.35, 0.5, icon, size=28, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, 2.7, 1.35, 0.4, f'Step {num}', size=11, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, 3.1, 1.35, 1.0, label, size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < len(steps) - 1:
        add_text(slide, x + 1.55, 2.8, 0.3, 0.5, '→', size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Data flow description
add_shape_bg(slide, 0.8, 5.0, 11.5, 1.8)
add_text(slide, 1.2, 5.2, 10.8, 0.4, 'Data Flow', size=18, bold=True, color=ACCENT)
add_text(slide, 1.2, 5.7, 10.8, 1.0,
    'User uploads resume → Django parses text (PyPDF2/python-docx) → Groq LLaMA extracts skills →\n'
    'Groq generates 10 technical questions → User answers via text or voice → Each answer evaluated\n'
    'by AI in real-time → Comprehensive report generated with scores and recommendation.',
    size=14, color=GRAY)

# ════════════════════════════════════════════
# SLIDE 6 — Key Features Deep Dive
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.8, '🌟  Key Features', size=36, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.2, 2)

features_detail = [
    ('🎯 Skill-Based Questions', 'Questions are generated specifically from skills found in the resume — not generic. Covers programming languages, frameworks, databases, and tools.'),
    ('🔊 Voice-Powered Interview', 'Text-to-Speech reads questions aloud. Speech-to-Text captures spoken answers via browser microphone. No extra APIs needed — uses Web Speech API.'),
    ('⚡ Real-Time Evaluation', 'Each answer is evaluated immediately by AI with a score (1-10) and detailed feedback, so the candidate can see how they did.'),
    ('📋 Comprehensive Report', 'Final report includes: overall score, per-skill breakdown, strengths, weaknesses, detailed feedback, and a hiring recommendation.'),
    ('💸 Completely Free', 'Uses Groq free tier API — no credit card, no paid subscription. Fast inference with LLaMA 3.3 70B model.'),
]

for i, (title, desc) in enumerate(features_detail):
    y = 1.5 + i * 1.1
    add_text(slide, 1.2, y, 4, 0.4, title, size=17, bold=True, color=WHITE)
    add_text(slide, 5.5, y, 7, 0.9, desc, size=14, color=GRAY)

# ════════════════════════════════════════════
# SLIDE 7 — Database Models
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.8, '🗃️  Database Design', size=36, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.2, 2.5)

# InterviewSession card
add_shape_bg(slide, 0.8, 1.6, 3.8, 4.5)
add_text(slide, 1.1, 1.8, 3.3, 0.4, 'InterviewSession', size=18, bold=True, color=ACCENT)
fields1 = ['resume_file (FileField)', 'resume_text (TextField)', 'skills_json (TextField)', 'status (CharField)', 'candidate_name (CharField)', 'report_text (TextField)', 'overall_score (FloatField)', 'created_at (DateTimeField)']
for i, f in enumerate(fields1):
    add_text(slide, 1.3, 2.4 + i * 0.42, 3.1, 0.4, f'• {f}', size=12, color=GRAY)

# Question card
add_shape_bg(slide, 5.0, 1.6, 3.5, 3.2)
add_text(slide, 5.3, 1.8, 3, 0.4, 'Question', size=18, bold=True, color=ACCENT)
fields2 = ['session (FK → Session)', 'question_text (TextField)', 'skill_category (CharField)', 'difficulty (CharField)', 'order (IntegerField)']
for i, f in enumerate(fields2):
    add_text(slide, 5.5, 2.4 + i * 0.42, 2.8, 0.4, f'• {f}', size=12, color=GRAY)

# Answer card
add_shape_bg(slide, 8.9, 1.6, 3.5, 3.0)
add_text(slide, 9.2, 1.8, 3, 0.4, 'Answer', size=18, bold=True, color=ACCENT)
fields3 = ['question (FK → Question)', 'answer_text (TextField)', 'evaluation (TextField)', 'score (IntegerField)', 'created_at (DateTimeField)']
for i, f in enumerate(fields3):
    add_text(slide, 9.4, 2.4 + i * 0.42, 2.8, 0.4, f'• {f}', size=12, color=GRAY)

# Relationship arrows
add_text(slide, 4.6, 2.8, 0.5, 0.4, '→', size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text(slide, 8.5, 2.8, 0.5, 0.4, '→', size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text(slide, 0.8, 6.4, 11, 0.5, 'Relationship: InterviewSession (1) → Questions (many) → Answers (many)', size=14, color=GRAY)

# ════════════════════════════════════════════
# SLIDE 8 — Project Structure
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.8, '📁  Project Structure', size=36, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.2, 2.5)

add_shape_bg(slide, 0.8, 1.6, 5.5, 5.2)

structure = [
    'bot/',
    '├── manage.py',
    '├── .env',
    '├── requirements.txt',
    '├── bot_project/',
    '│   ├── settings.py',
    '│   ├── urls.py',
    '│   └── wsgi.py',
    '├── interview/',
    '│   ├── models.py',
    '│   ├── views.py',
    '│   ├── gemini_service.py  (Groq AI)',
    '│   ├── resume_parser.py',
    '│   ├── forms.py',
    '│   ├── urls.py',
    '│   ├── templates/interview/',
    '│   │   ├── base.html',
    '│   │   ├── upload.html',
    '│   │   ├── interview.html',
    '│   │   └── report.html',
    '│   └── static/interview/css/',
    '└── media/resumes/',
]

for i, line in enumerate(structure):
    color = ACCENT if any(x in line for x in ['bot/', 'bot_project/', 'interview/', 'templates/', 'static/']) else GRAY
    add_text(slide, 1.1, 1.8 + i * 0.23, 5, 0.25, line, size=11, color=color, font_name='Consolas')

# File descriptions
add_shape_bg(slide, 6.8, 1.6, 5.5, 5.2)
add_text(slide, 7.1, 1.8, 5, 0.4, 'Key Files', size=18, bold=True, color=ACCENT)
file_descs = [
    ('gemini_service.py', 'All AI interactions — skill\nextraction, question gen,\nanswer evaluation, reports'),
    ('views.py', 'Request handling — upload,\nQ&A API endpoints, report\ngeneration'),
    ('resume_parser.py', 'PDF & DOCX text extraction\nusing PyPDF2 and python-docx'),
    ('interview.html', 'AJAX Q&A page with voice\ninteraction (TTS + STT)'),
    ('report.html', 'Animated report with SVG\nscore ring and skill bars'),
]

for i, (fname, desc) in enumerate(file_descs):
    y = 2.3 + i * 0.95
    add_text(slide, 7.3, y, 2.2, 0.3, fname, size=13, bold=True, color=WHITE)
    add_text(slide, 9.5, y, 3, 0.8, desc, size=11, color=GRAY)

# ════════════════════════════════════════════
# SLIDE 9 — UI Screenshots description
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.8, '🎨  User Interface', size=36, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.2, 2)

# Three UI screens
screens = [
    ('Upload Page', '• Drag-and-drop resume upload\n• File format validation\n• Name input field\n• Recent sessions list\n• Processing animation'),
    ('Interview Page', '• One question at a time\n• Progress bar tracking\n• Skill & difficulty tags\n• 🔊 Listen button (TTS)\n• 🎤 Mic button (STT)\n• Real-time score feedback'),
    ('Report Page', '• Animated SVG score ring\n• Strengths & weaknesses\n• Per-skill performance bars\n• Q&A breakdown with scores\n• Hiring recommendation\n• Print-friendly layout'),
]

for i, (title, features) in enumerate(screens):
    x = 0.8 + i * 4.1
    add_shape_bg(slide, x, 1.6, 3.8, 5.0)
    add_text(slide, x + 0.3, 1.8, 3.2, 0.5, title, size=20, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + 0.3, 2.5, 3.2, 4.0, features, size=14, color=GRAY)

add_text(slide, 0.8, 6.8, 11, 0.4, 'Design: Dark glassmorphism theme • Gradient accents • Smooth animations • Fully responsive', size=14, color=ACCENT_LIGHT)

# ════════════════════════════════════════════
# SLIDE 10 — How Groq AI Works
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.8, '🤖  AI Integration (Groq)', size=36, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.2, 2.5)

ai_functions = [
    ('extract_skills( )', 'Analyzes resume text and returns a JSON array of all technical skills found — programming languages, frameworks, tools, databases.', GREEN),
    ('generate_questions( )', 'Creates 10 tailored interview questions from the extracted skills with a mix of Easy, Medium, and Hard difficulty levels.', ACCENT),
    ('evaluate_answer( )', 'Scores each answer 1-10 with detailed feedback explaining what was good and what was missed.', YELLOW),
    ('generate_report( )', 'Produces a comprehensive JSON report with overall score, per-skill scores, strengths, weaknesses, and hiring recommendation.', RED),
]

for i, (func, desc, color) in enumerate(ai_functions):
    y = 1.6 + i * 1.35
    add_shape_bg(slide, 0.8, y, 11.5, 1.15)
    add_text(slide, 1.2, y + 0.15, 3.5, 0.4, func, size=16, bold=True, color=color, font_name='Consolas')
    add_text(slide, 5.0, y + 0.15, 7, 0.8, desc, size=14, color=GRAY)

add_text(slide, 0.8, 7.1, 11, 0.3, 'All AI calls use structured JSON prompts → LLaMA 3.3 70B processes → JSON response parsed with fallbacks', size=13, color=ACCENT_LIGHT)

# ════════════════════════════════════════════
# SLIDE 11 — Future Scope
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.8, '🚀  Future Scope', size=36, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.2, 2)

future_items = [
    ('👤  User Authentication', 'Add login/signup so candidates can track interview history'),
    ('📹  Video Recording', 'Record video responses for behavioral interview analysis'),
    ('🏢  Multi-Role Support', 'Customize questions for different job roles and levels'),
    ('📧  Email Reports', 'Auto-send PDF reports to candidates and recruiters'),
    ('📈  Analytics Dashboard', 'Admin panel with aggregate stats across all interviews'),
    ('🌐  Multi-Language', 'Support interviews in Hindi, Spanish, and other languages'),
]

for i, (title, desc) in enumerate(future_items):
    y = 1.5 + i * 0.9
    col = 0.8 if i < 3 else 7.0
    row_y = y if i < 3 else 1.5 + (i - 3) * 0.9
    add_shape_bg(slide, col, row_y, 5.5, 0.7)
    add_text(slide, col + 0.3, row_y + 0.05, 2.5, 0.4, title, size=15, bold=True, color=WHITE)
    add_text(slide, col + 0.3, row_y + 0.4, 5, 0.3, desc, size=12, color=GRAY)

# ════════════════════════════════════════════
# SLIDE 12 — Thank You
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(-2), Inches(8), Inches(8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(30, 27, 75)
shape.line.fill.background()

add_text(slide, 1, 2.0, 11, 1.2, 'Thank You!', size=56, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, 5.5, 3.5, 2.5)
add_text(slide, 1, 3.8, 11, 0.8, 'AI-Powered Technical Interview Bot', size=24, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 4.8, 11, 0.6, 'Built with Django  •  Groq AI  •  LLaMA 3.3 70B', size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 5.8, 11, 0.5, 'Questions & Feedback Welcome 🙌', size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# ─── Save ───
output_path = 'AI_Interview_Bot_Presentation.pptx'
prs.save(output_path)
print(f'✅ Presentation saved: {output_path}')
print(f'📄 Total slides: {len(prs.slides)}')
